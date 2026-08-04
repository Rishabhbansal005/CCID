import io
from typing import Dict, Any, List, Optional
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime

def _add_title_slide(prs, case: Dict[str, Any]):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = f"Investigation Report\n{case.get('title', 'Untitled')}"
    
    date_str = datetime.utcnow().strftime('%Y-%m-%d UTC')
    case_number = case.get('case_number', 'N/A')
    risk = case.get('priority', 'Unknown').upper()
    subtitle.text = f"Case Number: {case_number}\nPriority: {risk}\nGenerated: {date_str}"


def _add_executive_summary(prs, findings: List[Dict[str, Any]], risk_assessment: Optional[Dict[str, Any]]):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Executive Summary"
    
    tf = body_shape.text_frame
    
    critical_count = sum(1 for f in findings if f.get("severity") == "critical")
    high_count = sum(1 for f in findings if f.get("severity") == "high")
    
    p = tf.paragraphs[0]
    p.text = f"Total Findings: {len(findings)}"
    
    p = tf.add_paragraph()
    p.text = f"Critical Findings: {critical_count}"
    
    p = tf.add_paragraph()
    p.text = f"High Findings: {high_count}"
    
    if risk_assessment:
        p = tf.add_paragraph()
        p.text = f"Overall Risk Score: {risk_assessment.get('overall_risk_score', 'N/A')} / 25"
        
        p = tf.add_paragraph()
        p.text = f"Risk Level: {risk_assessment.get('risk_level', 'Unknown').upper()}"


def _add_findings_slide(prs, findings: List[Dict[str, Any]]):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Key Findings"
    
    tf = body_shape.text_frame
    tf.clear()
    
    # Sort by severity conceptually (critical first), but let's just show top 5
    priority_order = {"critical": 0, "high": 1, "medium": 2, "low": 3}
    sorted_findings = sorted(findings, key=lambda x: priority_order.get(x.get("severity", "low"), 4))
    
    if not sorted_findings:
        p = tf.paragraphs[0]
        p.text = "No findings recorded."
        return

    for i, finding in enumerate(sorted_findings[:5]):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        severity = finding.get('severity', 'low').upper()
        p.text = f"[{severity}] {finding.get('title', 'Untitled')}"
        p.font.size = Pt(18)
        
        # Optionally add a sub-bullet for description if it fits
        desc = finding.get('description', '')[:100] + ('...' if len(finding.get('description', '')) > 100 else '')
        if desc:
            p2 = tf.add_paragraph()
            p2.text = desc
            p2.font.size = Pt(14)
            p2.level = 1


def _add_evidence_summary(prs, evidence_list: List[Dict[str, Any]]):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = "Evidence Inventory"
    tf = body_shape.text_frame
    tf.clear()
    
    p = tf.paragraphs[0]
    p.text = f"Total Evidence Items Collected: {len(evidence_list)}"
    
    if not evidence_list:
        return
        
    for i, ev in enumerate(evidence_list[:5]):
        p2 = tf.add_paragraph()
        p2.text = f"{ev.get('evidence_number', 'EV-?')} : {ev.get('original_file_name', 'Unknown')} ({ev.get('evidence_type', 'digital')})"
        p2.level = 1


def generate_case_ppt(
    case: Dict[str, Any],
    evidence_list: List[Dict[str, Any]],
    findings: List[Dict[str, Any]],
    timeline_events: List[Dict[str, Any]],
    risk_assessment: Optional[Dict[str, Any]],
    investigator: Dict[str, Any],
    config: Dict[str, Any]
) -> bytes:
    """Generates a PowerPoint presentation based on the case data."""
    prs = Presentation()
    
    # Slide 1: Title
    _add_title_slide(prs, case)
    
    # Slide 2: Executive Summary
    if config.get("include_executive_summary", True):
        _add_executive_summary(prs, findings, risk_assessment)
        
    # Slide 3: Evidence Summary
    if config.get("include_evidence_list", True):
        _add_evidence_summary(prs, evidence_list)
        
    # Slide 4: Key Findings
    if config.get("include_findings", True):
        _add_findings_slide(prs, findings)
        
    # Save to BytesIO
    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    
    return ppt_stream.read()
